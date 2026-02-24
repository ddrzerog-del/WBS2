import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import re
import matplotlib.pyplot as plt
import matplotlib.patches as patches

# --- 1. 데이터 파싱 및 트리 구조화 ---
def parse_data(df):
    structured_data = []
    for _, row in df.iterrows():
        # A열(코드), B열(명칭) 가져오기
        code = str(row[0]).strip() if pd.notnull(row[0]) else ""
        name = str(row[1]).strip() if pd.notnull(row[1]) else ""
        
        if not code: continue

        # 코드에서 레벨 계산 (예: 1.1.2 -> 레벨 3)
        match = re.match(r'^([\d\.]+)', code)
        if match:
            clean_code = match.group(1).rstrip('.')
            level = clean_code.count('.') + 1
            structured_data.append({
                'id_code': clean_code,
                'code_text': code,
                'name_text': name,
                'level': level
            })
    return structured_data

def build_tree(data):
    # 숫자 순서대로 정렬 (1.2 다음 1.10이 오도록 처리)
    def split_code(c):
        return [int(i) for i in c['id_code'].split('.')]
    
    data.sort(key=split_code)
    
    nodes = {}
    root_nodes = []
    for item in data:
        code = item['id_code']
        node = {
            'code': code, 
            'code_text': item['code_text'], 
            'name_text': item['name_text'], 
            'level': item['level'], 
            'children': []
        }
        nodes[code] = node
        parts = code.split('.')
        if len(parts) > 1:
            parent_code = ".".join(parts[:-1])
            if parent_code in nodes:
                nodes[parent_code]['children'].append(node)
            else:
                if item['level'] == 1: root_nodes.append(node)
        else:
            root_nodes.append(node)
    return root_nodes

# --- 2. 좌표 및 레이아웃 계산 ---
def calculate_layout(root_nodes, config):
    layout_data = []
    wbs_w = config['wbs_w']
    wbs_h = config['wbs_h']
    l1_gap_x = config['l1_gap_x']
    l2_gap_x = config['l2_gap_x']
    v_gap_a = config['v_gap_a']
    extra_gaps = {
        3: config['extra_l3'],
        4: config['extra_l4'],
        5: config['extra_l5']
    }

    start_x = (33.8 - wbs_w) / 2
    start_y = (19.05 - wbs_h) / 2

    l1_count = len(root_nodes)
    if l1_count == 0: return []
    l1_width = (wbs_w - (l1_gap_x * (l1_count - 1))) / l1_count

    for i, l1 in enumerate(root_nodes):
        x_l1 = start_x + (i * (l1_width + l1_gap_x))
        y_l1 = start_y
        h_l1 = 1.0
        layout_data.append({'node': l1, 'x': x_l1, 'y': y_l1, 'w': l1_width, 'h': h_l1, 'level': 1})

        if l1['children']:
            l2_count = len(l1['children'])
            l2_width = (l1_width - (l2_gap_x * (l2_count - 1))) / l2_count
            current_y_for_l2 = y_l1 + h_l1 + v_gap_a

            for j, l2 in enumerate(l1['children']):
                x_l2 = x_l1 + (j * (l2_width + l2_gap_x))
                y_l2 = current_y_for_l2
                h_l2 = 1.0
                layout_data.append({'node': l2, 'x': x_l2, 'y': y_l2, 'w': l2_width, 'h': h_l2, 'level': 2})

                def stack_recursive(parent_node, px, py, pw, ph):
                    nonlocal layout_data
                    last_y = py + ph
                    for idx, child in enumerate(parent_node['children']):
                        lvl = child['level']
                        gap = v_gap_a + (extra_gaps.get(lvl, 0) if idx > 0 else 0)
                        target_y = last_y + gap
                        reduction = 0.3 * (lvl - 2)
                        c_w = max(l2_width - reduction, 4.0) 
                        c_x = (px + pw) - c_w
                        c_h = 0.8
                        layout_data.append({'node': child, 'x': c_x, 'y': target_y, 'w': c_w, 'h': c_h, 'level': lvl})
                        last_y = stack_recursive(child, c_x, target_y, c_w, c_h)
                    return last_y

                stack_recursive(l2, x_l2, y_l2, l2_width, h_l2)
    return layout_data

# --- 3. 실시간 미리보기 (Matplotlib) ---
def draw_preview(layout_data, code_width_cm):
    fig, ax = plt.subplots(figsize=(12, 6.75))
    ax.set_xlim(0, 33.8)
    ax.set_ylim(0, 19.05)
    ax.invert_yaxis()
    
    # 슬라이드 배경
    ax.add_patch(patches.Rectangle((0, 0), 33.8, 19.05, linewidth=1, edgecolor='#cccccc', facecolor='#ffffff'))

    for item in layout_data:
        x, y, w, h = item['x'], item['y'], item['w'], item['h']
        lvl = item['level']
        node = item['node']

        if lvl == 1: color, text_c = '#1f497d', 'white'
        elif lvl == 2: color, text_c = '#365f91', 'white'
        else: color, text_c = '#f2f2f2', 'black'

        # 1. 메인 박스
        ax.add_patch(patches.Rectangle((x, y), w, h, linewidth=0.5, edgecolor='#aaaaaa', facecolor=color))

        # 2. 코드 박스
        code_color = '#14325a' if lvl <= 2 else '#dddddd'
        code_text_c = 'white' if lvl <= 2 else 'black'
        ax.add_patch(patches.Rectangle((x, y), code_width_cm, h, linewidth=0.5, edgecolor='#666666', facecolor=code_color))

        # 3. 텍스트
        ax.text(x + code_width_cm/2, y + h/2, node['code_text'], color=code_text_c, fontsize=6, ha='center', va='center', fontweight='bold')
        display_name = (node['name_text'][:20] + '..') if len(node['name_text']) > 20 else node['name_text']
        ax.text(x + code_width_cm + 0.2, y + h/2, display_name, color=text_c, fontsize=6, ha='left', va='center')

    ax.set_axis_off()
    st.pyplot(fig)

# --- 4. PPT 생성 (TypeError 수정 버전) ---
def generate_ppt(layout_data, code_width_cm):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(33.8), Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for item in layout_data:
        node = item['node']
        lvl = item['level']
        x, y, w, h = item['x'], item['y'], item['w'], item['h']

        # 1. 메인 박스 (명칭용)
        main_shp = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
        main_shp.line.color.rgb = RGBColor(150, 150, 150)
        
        # [수정] 단색 채우기 선언 후 색상 적용
        main_shp.fill.solid()
        if lvl == 1:
            main_shp.fill.fore_color.rgb = RGBColor(31, 73, 125)
        elif lvl == 2:
            main_shp.fill.fore_color.rgb = RGBColor(54, 95, 145)
        else:
            c = min(230 + (lvl * 3), 250)
            main_shp.fill.fore_color.rgb = RGBColor(c, c, c+5)

        # 메인 텍스트 (명칭)
        tf = main_shp.text_frame
        tf.text = node['name_text']
        tf.margin_left = Cm(code_width_cm + 0.2) # 코드박스만큼 여백
        tf.vertical_anchor = 1 # Middle 정렬
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(9 if lvl > 2 else 11)
        p.font.color.rgb = RGBColor(255, 255, 255) if lvl <= 2 else RGBColor(0, 0, 0)

        # 2. 코드 박스 (A열용 - 메인박스 위에 덮음)
        code_shp = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(code_width_cm), Cm(h))
        
        # [수정] 단색 채우기 선언 후 색상 적용
        code_shp.fill.solid()
        if lvl <= 2:
            code_shp.fill.fore_color.rgb = RGBColor(20, 50, 90)
        else:
            code_shp.fill.fore_color.rgb = RGBColor(200, 200, 200)
        
        code_tf = code_shp.text_frame
        code_tf.text = node['code_text']
        code_p = code_tf.paragraphs[0]
        code_p.alignment = PP_ALIGN.CENTER
        code_p.font.size = Pt(8 if lvl > 2 else 10)
        code_p.font.bold = True
        code_p.font.color.rgb = RGBColor(255, 255, 255) if lvl <= 2 else RGBColor(0, 0, 0)

    return prs

# --- 5. Streamlit UI ---
st.set_page_config(page_title="WBS Master Pro", layout="wide")

st.sidebar.title("🎨 디자인 및 간격 설정")
code_w_input = st.sidebar.slider("코드 박스 너비 (cm)", 1.0, 5.0, 2.5, 0.1)

with st.sidebar.expander("📏 캔버스 및 기본 간격", expanded=True):
    wbs_w = st.number_input("WBS 전체 너비", 10.0, 32.0, 31.0)
    wbs_h = st.number_input("WBS 전체 높이", 5.0, 18.0, 16.0)
    v_gap_a = st.number_input("기준 수직 간격 (A)", 0.0, 5.0, 0.4, 0.05)
    l1_gap_x = st.number_input("L1 좌우 간격", 0.0, 10.0, 1.2)
    l2_gap_x = st.number_input("L2 좌우 간격", 0.0, 5.0, 0.4)

with st.sidebar.expander("↕️ 그룹간 추가 여백 (줄기 변경 시)"):
    extra_l3 = st.number_input("L3 그룹 간 추가", 0.0, 5.0, 0.3)
    extra_l4 = st.number_input("L4 그룹 간 추가", 0.0, 5.0, 0.2)
    extra_l5 = st.number_input("L5+ 그룹 간 추가", 0.0, 5.0, 0.1)

config = {
    'wbs_w': wbs_w, 'wbs_h': wbs_h, 'l1_gap_x': l1_gap_x, 'l2_gap_x': l2_gap_x,
    'v_gap_a': v_gap_a, 'extra_l3': extra_l3, 'extra_l4': extra_l4, 'extra_l5': extra_l5
}

st.title("📊 WBS 마스터 디자이너")
st.info("💡 엑셀의 A열은 코드(예: 1.1), B열은 작업명칭으로 구성해주세요.")

uploaded_file = st.file_uploader("엑셀 파일(.xlsx) 업로드", type=["xlsx"])

if uploaded_file:
    df = pd.read_excel(uploaded_file)
    raw_data = parse_data(df)
    
    if raw_data:
        tree = build_tree(raw_data)
        layout_data = calculate_layout(tree, config)
        
        st.subheader("🖼️ 슬라이드 디자인 미리보기")
        draw_preview(layout_data, code_w_input)
        
        if st.button("🚀 PPT 생성 및 다운로드", use_container_width=True):
            try:
                final_ppt = generate_ppt(layout_data, code_w_input)
                ppt_io = io.BytesIO()
                final_ppt.save(ppt_io)
                ppt_io.seek(0)
                st.download_button("🎁 PPT 파일 다운로드", ppt_io, "Smart_WBS_Layout.pptx")
            except Exception as e:
                st.error(f"PPT 생성 중 오류 발생: {e}")
